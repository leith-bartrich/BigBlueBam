"""
Generate BigBlueBam overview PowerPoint deck.
Usage: python scripts/build-deck.py
Output: docs/BigBlueBam-Overview.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
IMG = os.path.join(ROOT, "images")
OUT = os.path.join(ROOT, "docs", "BigBlueBam-Overview.pptx")

# ── Brand colours ────────────────────────────────────────────────────
BLUE = RGBColor(0x25, 0x63, 0xEB)  # primary-600
DARK = RGBColor(0x18, 0x18, 0x1B)  # zinc-950
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ZINC100 = RGBColor(0xF4, 0xF4, 0xF5)
ZINC400 = RGBColor(0xA1, 0xA1, 0xAA)
ZINC600 = RGBColor(0x52, 0x52, 0x5B)
ZINC800 = RGBColor(0x27, 0x27, 0x2A)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xD9, 0x77, 0x06)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


# ── Helpers ──────────────────────────────────────────────────────────

def _img(name):
    return os.path.join(IMG, name)


def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return p


def add_para(tf, text, size=18, color=WHITE, bold=False, space_before=Pt(6), font_name="Segoe UI"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    if space_before:
        p.space_before = space_before
    return p


def add_bullet(tf, text, size=16, color=WHITE, level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = "Segoe UI"
    p.level = level
    p.space_before = Pt(4)
    return p


def blue_chip(slide, left, top, label):
    """Small rounded-rect badge."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.4), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    set_text(tf, label, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    return shape


def section_divider(title, subtitle="", color=BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    # Title
    tb = text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1))
    set_text(tb.text_frame, title, size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        tb2 = text_box(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.8))
        set_text(tb2.text_frame, subtitle, size=22, color=ZINC400, align=PP_ALIGN.CENTER)
    return slide


def screenshot_slide(title, img_file, caption="", title_color=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    # Title
    tb = text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6))
    set_text(tb.text_frame, title, size=28, color=title_color, bold=True)
    # Image
    img_path = _img(img_file)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.8), Inches(1.1), Inches(11.7))
    # Caption
    if caption:
        tb2 = text_box(slide, Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.4))
        set_text(tb2.text_frame, caption, size=13, color=ZINC400, align=PP_ALIGN.CENTER)
    return slide


def two_screenshot_slide(title, img1, cap1, img2, cap2):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    tb = text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6))
    set_text(tb.text_frame, title, size=28, color=WHITE, bold=True)
    w = Inches(5.8)
    for i, (img_file, cap) in enumerate([(img1, cap1), (img2, cap2)]):
        x = Inches(0.5) + i * Inches(6.2)
        img_path = _img(img_file)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, x, Inches(1.2), w)
        ctb = text_box(slide, x, Inches(6.9), w, Inches(0.4))
        set_text(ctb.text_frame, cap, size=12, color=ZINC400, align=PP_ALIGN.CENTER)
    return slide


# =====================================================================
# SLIDE 1 — TITLE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

# Logo square
logo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.8), Inches(1.4), Inches(1.0), Inches(1.0))
logo.fill.solid()
logo.fill.fore_color.rgb = BLUE
logo.line.fill.background()
tf = logo.text_frame
tf.word_wrap = False
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
run = tf.paragraphs[0].add_run()
run.text = "B"
run.font.size = Pt(48)
run.font.bold = True
run.font.color.rgb = WHITE
run.font.name = "Segoe UI"

# Title text
tb = text_box(slide, Inches(1), Inches(2.8), Inches(11.3), Inches(1.2))
set_text(tb.text_frame, "BigBlueBam", size=56, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Tagline
tb2 = text_box(slide, Inches(1.5), Inches(4.1), Inches(10.3), Inches(1.0))
set_text(tb2.text_frame, "Project management built for human-AI teams.", size=26, color=ZINC400, align=PP_ALIGN.CENTER)
add_para(tb2.text_frame, "Engineers set the strategy. AI agents handle the grunt work. Everyone sees it on the board.", size=18, color=ZINC600, space_before=Pt(12))
tb2.text_frame.paragraphs[1].alignment = PP_ALIGN.CENTER

# Chips
for i, (label, c) in enumerate([
    ("140 MCP Tools", BLUE), ("14 Docker Services", PURPLE),
    ("530+ Tests", GREEN), ("4 Apps, 1 Stack", AMBER),
]):
    x = Inches(3.2) + i * Inches(1.85)
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.6), Inches(1.7), Inches(0.38))
    chip.fill.solid()
    chip.fill.fore_color.rgb = c
    chip.line.fill.background()
    ctf = chip.text_frame
    ctf.word_wrap = False
    set_text(ctf, label, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 2 — FOUR APPS, ONE PLATFORM
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
tb = text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7))
set_text(tb.text_frame, "Four apps, one platform", size=36, color=WHITE, bold=True)
add_para(tb.text_frame, "Shared auth, shared database, shared MCP server — zero integration overhead.", size=18, color=ZINC400, space_before=Pt(8))

apps = [
    ("Bam", "Kanban boards, sprints, tasks,\nproject dashboard, team management", BLUE, "03-board.png"),
    ("Helpdesk", "Customer support tickets,\nAI triage, auto-linked tasks", GREEN, "16-helpdesk-tickets.png"),
    ("Beacon", "Knowledge base, semantic search,\ngraph explorer, expiry governance", AMBER, "beacon-home.png"),
    ("Banter", "Team messaging, channels, DMs,\nthreads, voice/video, AI agents", PURPLE, "banter-channels.png"),
]

for i, (name, desc, color, img_file) in enumerate(apps):
    x = Inches(0.4) + i * Inches(3.2)
    y = Inches(1.8)
    # Color header bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.0), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    # App name
    ntb = text_box(slide, x, y + Inches(0.15), Inches(3.0), Inches(0.5))
    set_text(ntb.text_frame, name, size=24, color=color, bold=True)
    # Description
    dtb = text_box(slide, x, y + Inches(0.65), Inches(3.0), Inches(0.8))
    tf = dtb.text_frame
    tf.word_wrap = True
    set_text(tf, desc, size=13, color=ZINC400)
    # Screenshot
    img_path = _img(img_file)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, x, y + Inches(1.6), Inches(3.0))


# =====================================================================
# RAPID-FIRE FEATURE OVERVIEW (1-2 slides per major feature)
# =====================================================================
section_divider("Rapid-Fire Feature Tour", "One slide per major capability")

# Board
screenshot_slide("Kanban Board — Drag-and-Drop with Spring Physics", "03-board.png",
                 "5 configurable phases  ·  WIP limits  ·  Priority, assignee, points, due date on every card")

# Swimlanes + Views
two_screenshot_slide("Swimlanes & Multiple Views",
                     "05-swimlanes.png", "Swimlanes — group by assignee, priority, or epic",
                     "06-list-view.png", "List view — sortable columns, inline editing")

# Timeline + Calendar
two_screenshot_slide("Timeline & Calendar",
                     "07-timeline.png", "Gantt-style timeline with today marker",
                     "08-calendar.png", "Monthly calendar with task due dates")

# Task Detail + Dashboard
two_screenshot_slide("Task Detail & Project Dashboard",
                     "04-task-detail.png", "Rich-text, subtasks, attachments, comments, activity",
                     "09-project-dashboard.png", "Sprint progress, priority breakdown, overdue tasks")

# Command Palette + My Work
two_screenshot_slide("Power User Tools",
                     "14-command-palette.png", "Command palette — Ctrl+K to search and navigate",
                     "10-my-work.png", "My Work — cross-project view of your assignments")

# People Management
two_screenshot_slide("Organization & People Management",
                     "people-list.png", "Searchable, filterable, bulk-selectable people list",
                     "people-detail-overview.png", "Per-user detail with identity, projects, access, activity")

# SuperUser
two_screenshot_slide("SuperUser Console",
                     "superuser-overview.png", "Cross-org platform overview",
                     "superuser-people-list.png", "Global user management — every user across every org")

# Helpdesk
two_screenshot_slide("Helpdesk — Customer Support Portal",
                     "16-helpdesk-tickets.png", "Customers submit and track support tickets",
                     "17-helpdesk-conversation.png", "Ticket detail with full conversation thread")

# Beacon
two_screenshot_slide("Beacon — Knowledge Base",
                     "beacon-home.png", "Knowledge Home with stats, recent activity, quick actions",
                     "beacon-list.png", "Browse beacons by status, project, and tags")

two_screenshot_slide("Beacon — Graph & Governance",
                     "beacon-graph.png", "Knowledge Graph — connections via links and tag affinity",
                     "beacon-dashboard.png", "Governance dashboard — freshness score, at-risk beacons")

# Banter
two_screenshot_slide("Banter — Team Messaging",
                     "banter-channels.png", "Real-time channels with rich message compose",
                     "banter-search.png", "Full-text search across messages and channels")

# MCP
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
tb = text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7))
set_text(tb.text_frame, "140 MCP Tools — AI Agents at Full Parity", size=32, color=WHITE, bold=True)

categories = [
    ("Bam", "64 tools", BLUE, [
        "Tasks, sprints, comments, labels",
        "Board positions, bulk operations",
        "Sprint reports, analytics",
        "Members, settings, views",
    ]),
    ("Banter", "47 tools", PURPLE, [
        "Channels, DMs, threads",
        "Messages, reactions, pins",
        "Voice calls, admin settings",
        "Retention policies, search",
    ]),
    ("Beacon", "29 tools", AMBER, [
        "CRUD, publish, retire, verify",
        "Semantic + graph search",
        "Tags, links, saved queries",
        "Governance policies, graph",
    ]),
]

for i, (name, count, color, items) in enumerate(categories):
    x = Inches(0.6) + i * Inches(4.2)
    # Header
    htb = text_box(slide, x, Inches(1.3), Inches(3.8), Inches(0.5))
    set_text(htb.text_frame, f"{name} — {count}", size=22, color=color, bold=True)
    # Items
    itb = text_box(slide, x, Inches(1.9), Inches(3.8), Inches(3.0))
    itb.text_frame.word_wrap = True
    for j, item in enumerate(items):
        if j == 0:
            set_text(itb.text_frame, f"•  {item}", size=15, color=ZINC400)
        else:
            add_bullet(itb.text_frame, f"•  {item}", size=15, color=ZINC400)

# Agent workflow description
atb = text_box(slide, Inches(0.6), Inches(4.2), Inches(12), Inches(2.8))
atf = atb.text_frame
atf.word_wrap = True
set_text(atf, "What AI Agents Can Do", size=20, color=WHITE, bold=True)
for line in [
    "Create and manage tasks — set priority, assignee, move cards across phases, add subtasks",
    "Run sprints — create sprints, assign tasks, start/complete sprints, generate reports",
    "Triage helpdesk tickets — auto-create tasks, adjust priority, assign, respond to customers",
    "Manage knowledge via Beacon — create, search, verify, link, manage policies",
    "Collaborate on Banter — post to channels, reply to threads, summarize conversations",
    "Generate reports — sprint velocity, task distribution, team workload, overdue items",
]:
    add_bullet(atf, f"•  {line}", size=14, color=ZINC400)


# =====================================================================
# DEEP DIVES
# =====================================================================
section_divider("Deep Dive: Bam", "Kanban boards, sprints, views, and team management")

# Board deep dive
screenshot_slide("The Board — Your Shared Workspace", "03-board.png",
                 "Cards flow left to right through configurable phases  ·  WIP limits prevent bottlenecks  ·  Motion spring physics")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
tb = text_box(slide, Inches(0.6), Inches(0.3), Inches(6), Inches(0.6))
set_text(tb.text_frame, "Board Features", size=32, color=WHITE, bold=True)
ftb = text_box(slide, Inches(0.6), Inches(1.1), Inches(5.8), Inches(5.5))
ftf = ftb.text_frame
ftf.word_wrap = True
features = [
    ("5 Configurable Phases", "Customize column names, WIP limits, and done-states per project"),
    ("Drag-and-Drop", "Motion spring physics — cards animate naturally between phases"),
    ("Swimlanes", "Group by assignee, priority, or epic with collapsible rows"),
    ("5 Views", "Board, List, Timeline, Calendar, Workload — switch without losing filters"),
    ("Sprint Integration", "Assign tasks to sprints, track velocity, auto carry-forward"),
    ("Real-time Updates", "WebSocket + Redis PubSub — changes broadcast instantly to all users"),
    ("Optimistic UI", "TanStack Query mutations update immediately, rollback on failure"),
    ("Keyboard Shortcuts", "Ctrl+K command palette, vim-style navigation"),
]
for j, (feat, desc) in enumerate(features):
    if j == 0:
        set_text(ftf, feat, size=16, color=BLUE, bold=True)
    else:
        add_para(ftf, feat, size=16, color=BLUE, bold=True, space_before=Pt(14))
    add_para(ftf, desc, size=13, color=ZINC400, space_before=Pt(2))

# Task detail
if os.path.exists(_img("13-board-light.png")):
    slide.shapes.add_picture(_img("13-board-light.png"), Inches(6.8), Inches(1.1), Inches(6.0))

# Views deep dive
screenshot_slide("Task Detail — Everything in One Drawer", "04-task-detail.png",
                 "Rich text  ·  Subtasks  ·  File attachments  ·  Comments with reactions  ·  Activity feed  ·  Custom fields")

# Sprint & Dashboard
two_screenshot_slide("Sprints & Analytics",
                     "09-project-dashboard.png", "Project dashboard — velocity, burndown, priority breakdown",
                     "10-my-work.png", "My Work — cross-project task list, grouped by project")

# People management deep dive
section_divider("Deep Dive: People & Access", "Organization management, API keys, SuperUser console")

two_screenshot_slide("People Management",
                     "people-list.png", "Searchable list with role/status filters and bulk actions",
                     "people-detail-access.png", "Access tab — API keys, sessions, password management")

two_screenshot_slide("SuperUser Platform Console",
                     "superuser-context-banner.png", "Context-switched into another org — red banner warns operators",
                     "superuser-people-sessions.png", "Active sessions with IP, device, and revoke controls")

# ── Helpdesk deep dive ──────────────────────────────────────────────
section_divider("Deep Dive: Helpdesk", "Customer-facing support portal with AI triage")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
tb = text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7))
set_text(tb.text_frame, "Helpdesk → Board Pipeline", size=32, color=WHITE, bold=True)

steps = [
    ("1", "Customer submits ticket", "Client reports issue via the helpdesk portal\nwith category, priority, and description.", GREEN),
    ("2", "AI agent triages", "A task is auto-created on the board.\nAI sets priority, assigns, and responds.", AMBER),
    ("3", "Board syncs ticket", "Moving the task through phases auto-updates\nthe ticket status. Clients see progress.", BLUE),
]
for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.5) + i * Inches(4.2)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(1.4), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    ctf = circle.text_frame
    set_text(ctf, num, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Title
    stb = text_box(slide, x + Inches(0.8), Inches(1.4), Inches(3.2), Inches(0.5))
    set_text(stb.text_frame, title, size=20, color=WHITE, bold=True)
    # Desc
    dtb = text_box(slide, x, Inches(2.2), Inches(3.8), Inches(1.2))
    dtb.text_frame.word_wrap = True
    set_text(dtb.text_frame, desc, size=14, color=ZINC400)

# Helpdesk screenshots
two_screenshot_slide("Helpdesk Portal",
                     "16-helpdesk-tickets.png", "Customer ticket list with status and priority",
                     "18-helpdesk-detail-conversation.png", "Threaded conversation — agent and client messages")


# ── Beacon deep dive ────────────────────────────────────────────────
section_divider("Deep Dive: Beacon", "Knowledge base with expiry governance and semantic search")

screenshot_slide("Knowledge Home", "beacon-home.png",
                 "4,396 beacons  ·  440 at risk  ·  1,087 recently updated  ·  Quick actions and recent activity")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
tb = text_box(slide, Inches(0.6), Inches(0.3), Inches(6), Inches(0.6))
set_text(tb.text_frame, "Beacon Key Features", size=32, color=WHITE, bold=True)
ftb = text_box(slide, Inches(0.6), Inches(1.1), Inches(5.8), Inches(5.5))
ftf = ftb.text_frame
ftf.word_wrap = True
feats = [
    ("Expiry Governance", "Every beacon has a shelf life. Stale knowledge is surfaced before it misleads."),
    ("Semantic + Graph Search", "Hybrid retrieval: vector similarity (Qdrant), graph expansion, full-text fallback."),
    ("Knowledge Graph Explorer", "Visualize connections via typed links and implicit tag-affinity edges."),
    ("Agent Verification", "AI agents verify and challenge beacons within confidence bounds."),
    ("Versioned Content", "Every edit creates a new version with full diff history."),
    ("Governance Policies", "Org-level defaults with project-level overrides for verification intervals."),
    ("Saved Queries", "Named search configurations — private, project, or org-scoped."),
]
for j, (feat, desc) in enumerate(feats):
    if j == 0:
        set_text(ftf, feat, size=16, color=AMBER, bold=True)
    else:
        add_para(ftf, feat, size=16, color=AMBER, bold=True, space_before=Pt(14))
    add_para(ftf, desc, size=13, color=ZINC400, space_before=Pt(2))

if os.path.exists(_img("beacon-detail.png")):
    slide.shapes.add_picture(_img("beacon-detail.png"), Inches(6.8), Inches(1.1), Inches(6.0))

screenshot_slide("Knowledge Graph Explorer", "beacon-graph.png",
                 "Hub nodes sized by authority  ·  Freshness-colored rings  ·  At-risk pulsing  ·  Implicit + explicit edges")

two_screenshot_slide("Browse & Search",
                     "beacon-list.png", "Browse by status, project, tags — infinite scroll pagination",
                     "beacon-search.png", "Semantic search with multi-signal retrieval pipeline")

# ── Banter deep dive ────────────────────────────────────────────────
section_divider("Deep Dive: Banter", "Real-time team messaging with AI agent participation")

two_screenshot_slide("Channels & Search",
                     "banter-channels.png", "Real-time channels with rich compose and member list",
                     "banter-search.png", "Full-text search across all messages and channels")

two_screenshot_slide("Admin & Discovery",
                     "banter-admin.png", "Admin panel — channel management and retention policies",
                     "banter-browse.png", "Browse and discover channels across the organization")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
tb = text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7))
set_text(tb.text_frame, "Why Not Just Use Slack?", size=32, color=WHITE, bold=True)
ftb = text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(5.0))
ftf = ftb.text_frame
ftf.word_wrap = True
set_text(ftf, "Because Banter shares authentication, database, and deep cross-linking with BigBlueBam.", size=20, color=ZINC400)
reasons = [
    "Mention BBB-247 in a channel → links directly to the task",
    "AI agent triages a helpdesk ticket → posts update to #support-triage",
    "Sprint reports shared to channels with one click",
    "No webhooks, no bridges, no sync lag",
    "AI agents participate natively — same MCP tools, same permissions",
    "47 dedicated MCP tools for messaging automation",
]
for r in reasons:
    add_bullet(ftf, f"•  {r}", size=16, color=ZINC400)


# =====================================================================
# TECH STACK DEEP DIVE
# =====================================================================
section_divider("Deep Dive: Tech Stack", "Architecture, infrastructure, and developer experience")

# Architecture diagram (text-based)
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
tb = text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6))
set_text(tb.text_frame, "Architecture Overview", size=32, color=WHITE, bold=True)

# Left column — frontend
col_data = [
    ("Frontend Layer", BLUE, Inches(0.4), [
        "React 19 SPAs (Bam, Helpdesk, Beacon, Banter)",
        "TailwindCSS v4  ·  Motion (Framer Motion v11+)",
        "TanStack Query v5  ·  Zustand  ·  dnd-kit",
        "Radix UI  ·  Tiptap rich text  ·  Zod validation",
    ]),
    ("API Layer", GREEN, Inches(4.6), [
        "Node.js 22 LTS  ·  Fastify v5",
        "Drizzle ORM  ·  Zod (shared with frontend)",
        "WebSocket + Redis PubSub (real-time)",
        "BullMQ worker (email, export, notifications)",
    ]),
    ("Data Layer", AMBER, Inches(8.8), [
        "PostgreSQL 16 (RLS, JSONB, partitioned logs)",
        "Redis 7 (sessions, cache, PubSub, queues)",
        "MinIO / S3 (file attachments)",
        "Qdrant (vector search for Beacon)",
    ]),
]

for title, color, x, items in col_data:
    # Header
    hbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.2), Inches(4.0), Inches(0.06))
    hbar.fill.solid()
    hbar.fill.fore_color.rgb = color
    hbar.line.fill.background()
    htb = text_box(slide, x, Inches(1.35), Inches(4.0), Inches(0.5))
    set_text(htb.text_frame, title, size=20, color=color, bold=True)
    # Items
    itb = text_box(slide, x, Inches(1.95), Inches(4.0), Inches(2.5))
    itb.text_frame.word_wrap = True
    for j, item in enumerate(items):
        if j == 0:
            set_text(itb.text_frame, f"•  {item}", size=13, color=ZINC400)
        else:
            add_bullet(itb.text_frame, f"•  {item}", size=13, color=ZINC400)

# Infrastructure row
itb = text_box(slide, Inches(0.4), Inches(4.6), Inches(12.5), Inches(0.5))
set_text(itb.text_frame, "Infrastructure", size=20, color=PURPLE, bold=True)
ibar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(4.5), Inches(12.5), Inches(0.06))
ibar.fill.solid()
ibar.fill.fore_color.rgb = PURPLE
ibar.line.fill.background()

infra_items = [
    "Docker Compose — 14 services, single docker compose up",
    "nginx reverse proxy — all SPAs and APIs behind one port",
    "Turborepo + pnpm workspaces monorepo",
    "Multi-stage Dockerfiles — dev and production modes",
    "LiveKit SFU for voice/video  ·  Python voice agent (STT/TTS)",
    "MCP Server — Streamable HTTP + SSE + stdio transports",
]
iftb = text_box(slide, Inches(0.4), Inches(5.2), Inches(12.5), Inches(2.0))
iftb.text_frame.word_wrap = True
for j, item in enumerate(infra_items):
    if j == 0:
        set_text(iftb.text_frame, f"•  {item}", size=13, color=ZINC400)
    else:
        add_bullet(iftb.text_frame, f"•  {item}", size=13, color=ZINC400)


# Key numbers
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
tb = text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.6))
set_text(tb.text_frame, "Key Numbers", size=32, color=WHITE, bold=True)

numbers = [
    ("14", "Docker\nServices", BLUE),
    ("140", "MCP\nTools", PURPLE),
    ("530+", "Test\nCases", GREEN),
    ("40+", "Database\nTables", AMBER),
    ("4", "React\nSPAs", BLUE),
    ("38", "API Route\nModules", GREEN),
]

for i, (num, label, color) in enumerate(numbers):
    col = i % 6
    x = Inches(0.6) + col * Inches(2.1)
    y = Inches(1.5)
    # Number
    ntb = text_box(slide, x, y, Inches(1.8), Inches(1.0))
    set_text(ntb.text_frame, num, size=52, color=color, bold=True, align=PP_ALIGN.CENTER)
    # Label
    ltb = text_box(slide, x, y + Inches(1.0), Inches(1.8), Inches(0.8))
    ltb.text_frame.word_wrap = True
    set_text(ltb.text_frame, label, size=15, color=ZINC400, align=PP_ALIGN.CENTER)

# Monorepo structure
mtb = text_box(slide, Inches(0.6), Inches(3.6), Inches(12), Inches(3.5))
mtf = mtb.text_frame
mtf.word_wrap = True
set_text(mtf, "Monorepo Structure", size=22, color=WHITE, bold=True)
add_para(mtf, "", size=6, color=DARK, space_before=Pt(4))

mono_items = [
    "apps/api/ — Bam Fastify REST API + WebSocket (23 route modules, ~63 source files)",
    "apps/frontend/ — Bam React SPA (~55 source files, 8 pages, command palette)",
    "apps/banter-api/ — Banter Fastify API + WebSocket (15 routes, 18 DB tables)",
    "apps/banter/ — Banter React SPA (14 components, 7 pages)",
    "apps/beacon-api/ — Beacon Fastify API (knowledge base, search, graph, policies)",
    "apps/beacon/ — Beacon React SPA (knowledge home, graph explorer, editor)",
    "apps/helpdesk-api/ + apps/helpdesk/ — Customer support portal",
    "apps/mcp-server/ — MCP protocol server (140 tools, 10+ resources, 8 prompts)",
    "apps/worker/ — BullMQ background jobs (email, notifications, export, retention)",
    "packages/shared/ — Zod schemas, TypeScript types, constants",
]
for item in mono_items:
    add_bullet(mtf, f"•  {item}", size=12, color=ZINC400)


# =====================================================================
# CLOSING SLIDE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

# Logo
logo2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.8), Inches(1.8), Inches(1.0), Inches(1.0))
logo2.fill.solid()
logo2.fill.fore_color.rgb = BLUE
logo2.line.fill.background()
ltf = logo2.text_frame
ltf.paragraphs[0].alignment = PP_ALIGN.CENTER
run = ltf.paragraphs[0].add_run()
run.text = "B"
run.font.size = Pt(48)
run.font.bold = True
run.font.color.rgb = WHITE
run.font.name = "Segoe UI"

tb = text_box(slide, Inches(1), Inches(3.2), Inches(11.3), Inches(1.0))
set_text(tb.text_frame, "BigBlueBam", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

tb2 = text_box(slide, Inches(2), Inches(4.3), Inches(9.3), Inches(0.8))
set_text(tb2.text_frame, "Project management built for human-AI teams.", size=24, color=ZINC400, align=PP_ALIGN.CENTER)

urls = text_box(slide, Inches(2), Inches(5.5), Inches(9.3), Inches(1.2))
utf = urls.text_frame
utf.word_wrap = True
set_text(utf, "localhost/b3/  ·  localhost/beacon/  ·  localhost/banter/  ·  localhost/helpdesk/", size=16, color=ZINC600, align=PP_ALIGN.CENTER)
add_para(utf, "docker compose up  — and you're running.", size=18, color=ZINC400, space_before=Pt(16))
utf.paragraphs[1].alignment = PP_ALIGN.CENTER

# =====================================================================
# SAVE
# =====================================================================
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
